import os
import io
import time
import json
import aiofiles
import subprocess
from PIL import Image
import pytesseract
from pptx import Presentation

from app.utils.file_handler import change_to_processed
from app.services.image_caption import describe_image
from app.extractors.pdf_extractor import ensure_dirs
# from app.core.logger import #app_logger  


async def is_legacy_ppt(file_path: str) -> bool:
    return file_path.lower().endswith(".ppt") and not file_path.lower().endswith(".pptx")


async def convert_ppt_to_pptx(ppt_path: str) -> str:
    output_dir = os.path.dirname(ppt_path)
    try:
        #app_logger.info(f"Converting legacy PPT to PPTX: {ppt_path}")
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pptx", "--outdir", output_dir, ppt_path
        ], check=True)
        new_path = ppt_path.replace(".ppt", ".pptx")
        if os.path.exists(new_path):
            #app_logger.info(f"Conversion successful: {new_path}")
            return new_path
        else:
            #app_logger.error("Conversion failed: .pptx file not found")
            return None
    except Exception as e:
        #app_logger.exception("Error converting PPT to PPTX")
        raise RuntimeError(f"Failed to convert PPT to PPTX: {e}")


async def extract_ppt_content(file_path):
    start = time.time()
    #app_logger.info(f"Starting PPT extraction: {file_path}")

    if await is_legacy_ppt(file_path):
        #app_logger.debug("Detected legacy .ppt file")
        converted = await convert_ppt_to_pptx(file_path)
        if not converted:
            #app_logger.error("PPT conversion failed.")
            return {"error": "Failed to convert .ppt to .pptx"}
        file_path = converted

    try:
        prs = Presentation(file_path)
    except Exception as e:
        #app_logger.exception("Failed to load PPTX file")
        return {"error": str(e)}

    filename = os.path.splitext(os.path.basename(file_path))[0]
    img_root, img_summary_dir, img_vision_dir = await ensure_dirs()
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        texts, ocr_texts, image_descriptions = [], [], []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            if shape.shape_type == 13:  # Picture
                try:
                    img_bytes, ext = shape.image.blob, shape.image.ext
                    img_name = f"{filename}_slide{i}_img.{ext}"
                    img_path = os.path.join(img_root, img_name)

                    async with aiofiles.open(img_path, "wb") as f:
                        await f.write(img_bytes)

                    pil_img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(pil_img).strip()

                    if ocr_text:
                        ocr_texts.append(ocr_text)
                        summary_file = os.path.join(img_summary_dir, f"{filename}_slide{i}_summary.json")
                        async with aiofiles.open(summary_file, "w", encoding="utf-8") as f:
                            await f.write(json.dumps({"ocr_text": ocr_text}, indent=2))
                    else:
                        description = await describe_image(img_path)
                        image_descriptions.append(description)
                        vision_file = os.path.join(img_vision_dir, f"{filename}_slide{i}_vision.json")
                        async with aiofiles.open(vision_file, "w", encoding="utf-8") as f:
                            await f.write(json.dumps({"description": description}, indent=2))
                except Exception as e:
                    #app_logger.exception(f"Error processing image on slide {i}")
                    print(f"Error processing image on slide {i}: {e}")

        slides.append({
            "slide_number": i,
            "text_blocks": texts,
            "img_summary_texts": ocr_texts,
            "img_vision_descriptions": image_descriptions
        })

    result = {
        "metadata": {
            "file_name": os.path.basename(file_path),
            "file_type": "pptx",
            "file_size": f"{os.path.getsize(file_path)/1024/1024:.2f} MB",
            "slide_count": len(prs.slides)
        },
        "slides": slides,
        "summary": "PPT processed",
        "total_time_taken": f"{time.time() - start:.2f} sec"
    }

    output_path = f"output/{filename}.json"
    try:
        #app_logger.debug(f"Writing PPT output JSON: {output_path}")
        async with aiofiles.open(output_path, "w", encoding="utf-8") as f:
            await f.write(json.dumps(result, indent=2, ensure_ascii=False))
    except Exception as e:
        #app_logger.exception("Failed to write JSON output file")
        raise IOError(f"Failed to write JSON output file: {e}")

    await change_to_processed(str(file_path), "PPT")
    #app_logger.info(f"PPT extraction completed: {file_path}")
    return result
